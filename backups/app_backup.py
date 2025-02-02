import os
import base64
from flask import Flask, request, redirect, url_for, flash, render_template
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from zipfile import BadZipFile

app = Flask(__name__)
app.secret_key = "secret"

# Bullet point styles for indentation levels
bullet_styles = {
    0: "\u2022",  # Level 1
    1: "\u25e6",  # Level 2
    2: "\u25aa",  # Level 3
    3: "\u25ab",  # Level 4
    4: "-",       # Level 5
    5: "\u2013",  # Level 6
    6: "\u2794",  # Level 7
    7: "\u2192",  # Level 8
}

# Function to handle image embedding
def embed_image_as_base64(image_obj, images_list):
    """Convert raw image blob to base64 data URI, append to images_list."""
    blob = getattr(image_obj, "blob", None)
    if not blob:
        return
    base64data = base64.b64encode(blob).decode("utf-8")
    mime_type = "image/" + image_obj.ext.lower()
    data_uri = f"data:{mime_type};base64,{base64data}"
    images_list.append(data_uri)

# Function to parse PowerPoint file
def parse_pptx(filepath):
    try:
        prs = Presentation(filepath)
    except (PackageNotFoundError, BadZipFile):
        flash("Uploaded file is not a valid PowerPoint or is corrupted.")
        return redirect(url_for("index"))

    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        title = None
        youtube_links = []
        images = []
        text_html = ""
        table_html = ""

        for shape in slide.shapes:
            # Handle title
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                title = shape.text

            # Handle text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    bullet_level = paragraph.level
                    runs_html = ""

                    for run in paragraph.runs:
                        run_text = run.text.replace("<", "&lt;").replace(">", "&gt;")
                        if hasattr(run, "bold") and run.bold:
                            run_text = f"<strong>{run_text}</strong>"
                        if hasattr(run, "italic") and run.italic:
                            run_text = f"<em>{run_text}</em>"
                        runs_html += run_text

                    if runs_html.strip():
                        bullet_symbol = bullet_styles.get(bullet_level, "\u2022")  # Default to "•"
                        text_html += f"<li style='margin-left:{20 * bullet_level}px; list-style-type: none;'>{bullet_symbol} {runs_html}</li>"

            # Handle images
            if hasattr(shape, "image") and shape.image:
                embed_image_as_base64(shape.image, images)

            # Handle tables
            if shape.has_table:
                table = shape.table
                table_html += "<table border='1' style='border-collapse: collapse;'>"
                for row in table.rows:
                    table_html += "<tr>"
                    for cell in row.cells:
                        table_html += f"<td>{cell.text}</td>"
                    table_html += "</tr>"
                table_html += "</table>"

        if not title:
            title = f"Slide {slide_num}"

        slides_data.append({
            "title": title,
            "slide_number": slide_num,
            "text_html": f"<ul>{text_html}</ul>" if text_html else "",
            "table_html": table_html if table_html else "",
            "images": images,
            "youtube_links": youtube_links,
        })

    return slides_data

# Route for the index page
@app.route("/")
def index():
    return render_template("index.html")

# Route to handle file upload
@app.route("/upload", methods=["POST"])
def upload_pptx():
    if "file" not in request.files:
        flash("No file part")
        return redirect(request.url)

    file = request.files["file"]
    if file.filename == "":
        flash("No selected file")
        return redirect(request.url)

    if file:
        filepath = os.path.join("uploads", file.filename)
        file.save(filepath)

        slides_data = parse_pptx(filepath)

        return render_template("results.html", slides_data=slides_data)

# Run the app
if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5001)

