from flask import Flask, render_template, request, redirect, url_for, flash, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from zipfile import BadZipFile
from docx import Document
import os

app = Flask(__name__)
app.secret_key = "your_secret_key"
UPLOAD_FOLDER = "uploads"
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER


# Bullet styles for different levels
bullet_styles = {
    0: "•",  # Level 1
    1: "○",  # Level 2
    2: "▪",  # Level 3
    3: "▫",  # Level 4
    4: "-",  # Level 5
    5: "–",  # Level 6
    6: "»",  # Level 7
    7: "→"   # Level 8
}


def embed_image_as_base64(image_obj, images_list):
    """Helper function to embed image as base64."""
    blob = image_obj.blob
    if not blob:
        return
    ext = image_obj.ext.lower()
    mime_type = f"image/{ext}"
    base64data = blob.decode("utf-8")
    data_uri = f"data:{mime_type};base64,{base64data}"
    images_list.append(data_uri)


def parse_pptx(filepath):
    try:
        prs = Presentation(filepath)
    except (PackageNotFoundError, BadZipFile):
        flash("Uploaded file is not a valid PowerPoint or is corrupted.")
        return redirect(url_for("index"))

    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        title = None
        youtube_links = []
        images = []
        text_html = ""
        table_html = ""

        for shape in slide.shapes:
            # Handle title
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                title = shape.text

            # Handle text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    bullet_level = paragraph.level
                    runs_html = ""
                    for run in paragraph.runs:
                        run_text = run.text.replace("<", "&lt;").replace(">", "&gt;")
                        if hasattr(run, "bold") and run.bold:
                            run_text = f"<strong>{run_text}</strong>"
                        if hasattr(run, "italic") and run.italic:
                            run_text = f"<em>{run_text}</em>"
                        runs_html += run_text

                    if runs_html.strip():
                        bullet_symbol = bullet_styles.get(bullet_level, "•")
                        text_html += f"<li style='margin-left:{20 * bullet_level}px; list-style-type: none;'>{bullet_symbol} {runs_html}</li>"

            # Handle images
            if hasattr(shape, "image") and shape.image:
                embed_image_as_base64(shape.image, images)

            # Handle grouped shapes
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                for subshape in shape.shapes:
                    if hasattr(subshape, "image") and subshape.image:
                        embed_image_as_base64(subshape.image, images)

            # Handle background images
            if hasattr(shape, "fill") and hasattr(shape.fill, "picture") and shape.fill.picture:
                embed_image_as_base64(shape.fill.picture, images)

            # Handle tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_html = "<tr>"
                    for cell in row.cells:
                        row_html += f"<td>{cell.text}</td>"
                    row_html += "</tr>"
                    table_html += row_html

        if not title:
            title = f"Slide {slide_num}"

        slides_data.append({
            "title": title,
            "slide_number": slide_num,
            "text_html": f"<ul>{text_html}</ul>" if text_html else "",
            "table_html": f"<table>{table_html}</table>" if table_html else "",
            "images": images,
            "youtube_links": youtube_links
        })

    return slides_data


def generate_word_doc(slides_data):
    """Generate a Word document from slide data."""
    doc = Document()
    for slide in slides_data:
        doc.add_heading(slide["title"], level=1)
        if slide["text_html"]:
            doc.add_paragraph(slide["text_html"], style="List Bullet")
        if slide["table_html"]:
            doc.add_paragraph(slide["table_html"])
        if slide["images"]:
            doc.add_paragraph("Images are embedded.")
    word_path = os.path.join(UPLOAD_FOLDER, "parsed_output.docx")
    doc.save(word_path)
    return word_path


@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        if "file" not in request.files:
            flash("No file part")
            return redirect(request.url)
        file = request.files["file"]
        if file.filename == "":
            flash("No selected file")
            return redirect(request.url)
        if file:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)

            slides_data = parse_pptx(filepath)
            if isinstance(slides_data, str):
                return slides_data  # Handles flash redirect case

            word_path = generate_word_doc(slides_data)
            return send_file(word_path, as_attachment=True, download_name="parsed_output.docx")

    return render_template("index.html")


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5001)

