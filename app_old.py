import os
import uuid
import base64
from flask import Flask, request, render_template, redirect, url_for, flash
from werkzeug.utils import secure_filename
from pptx import Presentation

os.makedirs("static/slide_images", exist_ok=True)

app = Flask(__name__)
app.secret_key = "YOUR_SECRET_KEY_HERE"  # Replace with something random/secure

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

ALLOWED_EXTENSIONS = {"pptx"}

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")

@app.route("/upload", methods=["POST"])
def upload_pptx():
    if "pptx_file" not in request.files:
        flash("No file part in the request.")
        return redirect(url_for("index"))

    file = request.files["pptx_file"]
    if file.filename == "":
        flash("No selected file.")
        return redirect(url_for("index"))

    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(filepath)

        slides_data = parse_pptx(filepath)
        return render_template("results.html", slides_data=slides_data)
    else:
        flash("Invalid file type. Please upload a .pptx file.")
        return redirect(url_for("index"))

def parse_pptx(filepath):
    prs = Presentation(filepath)
    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        title = None
        youtube_links = []
        images = []
        text_html = ""

        for shape in slide.shapes:
            # If shape is a Title placeholder
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                title = shape.text

            # If shape has text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    bullet_level = paragraph.level
                    bullet_prefix = ""

                    runs_html = ""
                    for run in paragraph.runs:
                        run_text = run.text.replace("<", "&lt;").replace(">", "&gt;")
			# We'll track if we're currently inside a <ul> list
			in_list = False

			for paragraph in shape.text_frame.paragraphs:
				 bullet_level = paragraph.level
				runs_html = ""

				# Build runs_html (with bold/italic checks) as you already do
				for run in paragraph.runs:
					run_text = run.text.replace("<", "&lt;").replace(">", "&gt;")
       	
					 if hasattr(run, "bold") and run.bold:
						run_text = f"<strong>{run_text}</strong>"
					if hasattr(run, "italic") and run.italic:
						run_text = f"<em>{run_text}</em>"

					runs_html += run_text

				# If there's actual text (avoid empty bullet lines)
				if runs_html.strip():
 					# If bullet_level >= 0, treat it as a bullet
					if bullet_level >= 0:
						# If not already in a <ul>, start one
						if not in_list:
							text_html += "<ul>"
							in_list = True
           	

						 # Add <li>
						text_html += f"<li>{runs_html}</li>"
					else:
						# Not a bullet -> close a <ul> if open
						if in_list:
							text_html += "</ul>"
							in_list = False
            
						# Then treat as a normal paragraph
						text_html += f"<p>{runs_html}</p>"
					else:
						# If runs_html is blank, maybe skip or close a list
						pass

				# After finishing the paragraphs loop, if we ended inside a <ul>, close it
 				if in_list:
 				text_html += "</ul>"
				in_list = False

                        if hasattr(run, "bold") and run.bold:
	                        run_text = f"<strong>{run_text}</strong>"
                        if hasattr(run, "italic") and run.italic:
        	                run_text = f"<em>{run_text}</em>"
                        runs_html += run_text


                     # If shape is a Picture
                     if shape.shape_type == 13:  # PICTURE
                    	 image = shape.image
                          if image:
                     	 	 # Convert the raw bytes to base64
                                 base64data = base64.b64encode(image.blob).decode("utf-8")
                                 # Guess the MIME type based on the extension
                                 if image.ext.lower() in ["jpg", "jpeg"]:
                                	 mime_type = "image/jpeg"
                                 elif image.ext.lower() == "png":
                                	 mime_type = "image/png"
                                 else:
                                	 # fallback, e.g. "image/bmp", "image/gif", etc.
                                	 mime_type = f"image/{image.ext}"

                                 # Construct the data URI
                                 data_uri = f"data:{mime_type};base64,{base64data}"

                                 # Add to the images list (so we can <img src="{{ data_uri }}"> in HTML)
                                 images.append(data_uri)
	
            # If there's a hyperlink (check for YouTube)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            link = run.hyperlink.address
                            if "youtube.com" in link or "youtu.be" in link:
                                youtube_links.append(link)

        if not title:
            title = f"Slide {slide_num}"

        slides_data.append({
            "title": title,
            "slide_number": slide_num,
            "text_html": text_html,
            "images": images,
            "youtube_links": youtube_links
        })

    return slides_data

if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5001)

