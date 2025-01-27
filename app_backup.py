from flask import Flask, render_template, request, redirect, url_for, flash, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.exc import PackageNotFoundError
from zipfile import BadZipFile
from docx import Document

import os

app = Flask(__name__)
app.secret_key = "b5f8c27a6c7a4e8d9f561c6277e739bc"

# Ensure uploads folder exists
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# Route to handle the index page and file upload
@app.route('/upload', methods=['POST'])
def upload_file():
    if request.method == "POST":
        # Check if a file is included in the request
        if "file" not in request.files:
            flash("No file part")
            return redirect(request.url)
        file = request.files["file"]
        if file.filename == "":
            flash("No selected file")
            return redirect(request.url)
        if file:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)

            try:
                slides_data = parse_pptx(filepath)
                output_file = convert_to_word(slides_data)
                return redirect(url_for("download_file", filename=output_file))
            except Exception as e:
                flash(f"An error occurred: {str(e)}")
                return redirect(request.url)
    return render_template("index.html")

# Function to parse PowerPoint file
def parse_pptx(filepath):
    try:
        prs = Presentation(filepath)
    except (PackageNotFoundError, BadZipFile):
        raise Exception("Invalid or corrupted PowerPoint file.")

    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        title = None
        text_html = ""
        table_html = ""

        for shape in slide.shapes:
            # Handle title
            if shape.is_placeholder and shape.placeholder_format.type == 1:
                title = shape.text

            # Handle text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text_html += f"<p>{paragraph.text}</p>"

            # Handle tables
            if shape.has_table:
                table = shape.table
                table_html += "<table border='1'>"
                for row in table.rows:
                    table_html += "<tr>"
                    for cell in row.cells:
                        table_html += f"<td>{cell.text}</td>"
                    table_html += "</tr>"
                table_html += "</table>"

        slides_data.append({
            "title": title or f"Slide {slide_num}",
            "text_html": text_html,
            "table_html": table_html,
        })

    return slides_data

# Function to convert slides data to Word
def convert_to_word(slides_data):
    output_path = "output.docx"
    doc = Document()

    for slide in slides_data:
        doc.add_heading(slide["title"], level=1)
        if slide["text_html"]:
            doc.add_paragraph(slide["text_html"])
        if slide["table_html"]:
            doc.add_paragraph(slide["table_html"])
        doc.add_paragraph("\n")

    doc.save(output_path)
    return output_path

# Route to handle file download
@app.route("/download/<filename>")
def download_file(filename):
    file_path = os.path.join(os.getcwd(), filename)
    return send_file(file_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5001)

